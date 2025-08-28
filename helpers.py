import smtplib
import logging
import pymssql
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from config import settings

import pymssql
import io
import math
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from datetime import datetime


def update_items_from_xml(db: pymssql.Connection, table_name: str, xml_string: str, user_id: int, item_name: str):
    try:
        with db.cursor(as_dict=True) as cursor:
            sql_command = "EXEC usp_bulk_update @tableName=%s, @xmlText=%s, @userID=%d"
            cursor.execute(sql_command, (table_name, xml_string, user_id))
            
            result = cursor.fetchone()
            affected_rows = result.get('AffectedRowCount', 0) if result else 0

        db.commit()
        
        if affected_rows > 0:
            message = f"Data processed for '{item_name}'. {affected_rows} row(s) were updated."
        else:
            message = "Operation successful, but no changes were made to the data."

        return {
            "status": "success",
            "message": message,
            "affected_rows": affected_rows
        }
    
    except pymssql.Error as ex:
        logging.error(f"Database Service Error in update_items_from_xml: {ex}")
        raise
    except Exception as e:
        logging.error(f"Generic Service Error in update_items_from_xml: {e}")
        raise

def fetch_data(db: pymssql.Connection, proc_name: str, params: tuple = ()):
    rows = []
    try:
        with db.cursor(as_dict=True) as cursor:
            cursor.callproc(proc_name, params)
            rows = cursor.fetchall()
    except pymssql.Error as ex:
        logging.error(f"Database Helper Error in fetch_data for '{proc_name}': {ex}")
        raise
    
    return rows

def send_email(to_email: str, subject: str, html_content: str):
    sender_email = settings.SENDER_EMAIL
    sender_password = settings.SENDER_PASSWORD

    message = MIMEMultipart("alternative")
    message["Subject"] = subject
    message["From"] = sender_email
    message["To"] = to_email
    message.attach(MIMEText(html_content, "html"))

    try:
        with smtplib.SMTP("smtp.office365.com", 587) as server:
            server.starttls()
            server.login(sender_email, sender_password)
            server.sendmail(
                sender_email, to_email, message.as_string()
            )
            logging.info(f"Email sent successfully to {to_email}.")
            return True
            
    except smtplib.SMTPAuthenticationError:
        logging.error("Failed to send email: SMTP Authentication failed. Check SENDER_EMAIL/SENDER_PASSWORD.")
        return False
    except Exception as e:
        logging.error(f"An unexpected error occurred while sending email: {e}")
        return False
    

# # --- Configuration Constants (Powerpoint presentation)---
# SLIDE_WIDTH = Inches(13.33)
# SLIDE_HEIGHT = Inches(7.5)
# LEFT_MARGIN = Inches(0.5)
# RIGHT_MARGIN = Inches(0.5)
# TOP_MARGIN = Inches(0.5)
# LEFT_COLUMN_WIDTH = Inches(6.0)
# RIGHT_COLUMN_WIDTH = Inches(6.0)
# COLUMN_SPACING = Inches(0.3)
# TITLE_FONT_SIZE = Pt(18)
# TITLE_FONT_BOLD = True
# TITLE_FONT_COLOR = RGBColor(0, 51, 102)
# SECTION_FONT_SIZE = Pt(18)
# SECTION_FONT_BOLD = True
# SECTION_FONT_COLOR = RGBColor(0, 51, 102)
# SUBSECTION_FONT_SIZE = Pt(18)
# SUBSECTION_FONT_BOLD = True
# SUBSECTION_FONT_COLOR = RGBColor(0, 51, 102)
# BULLET_FONT_SIZE = Pt(12)
# BULLET_FONT_COLOR = RGBColor(64, 64, 64)
# TABLE_HEADER_FONT_SIZE = Pt(12)
# TABLE_HEADER_FONT_BOLD = True
# TABLE_HEADER_FONT_COLOR = RGBColor(255, 255, 255)
# TABLE_HEADER_BG_COLOR = RGBColor(0, 51, 102)
# TABLE_CELL_FONT_SIZE = Pt(12)
# TABLE_CELL_FONT_COLOR = RGBColor(0, 0, 0)
# SECTION_SPACING = Inches(0.2)
# MIN_ROW_HEIGHT = Inches(0.3)
# SUBSECTION_SPACING = Inches(0.15)
# BULLET_POINT_SPACING = Inches(0.0)


# def create_custom_presentation_from_xml(xml_string: str) -> io.BytesIO:
#     try:
#         root = ET.fromstring(xml_string)
#         prs = Presentation()
#         prs.slide_width = SLIDE_WIDTH
#         prs.slide_height = SLIDE_HEIGHT
#         blank_slide_layout = prs.slide_layouts[6]

#         for result in root.findall('Result'):
#             result_name = result.get('name')
            
#             # Process Page 1
#             page1 = result.find(".//Page[@number='1']")
#             if page1 is not None:
#                 slide = prs.slides.add_slide(blank_slide_layout)
#                 title_text = page1.find('Title').text if page1.find('Title') is not None and page1.find('Title').text else result_name
#                 add_slide_title(slide, title_text)
                
#                 current_y = TOP_MARGIN + Inches(0.4)
#                 left_x = LEFT_MARGIN
#                 right_x = left_x + LEFT_COLUMN_WIDTH + COLUMN_SPACING

#                 # Left Column Processing
#                 left_y = current_y
#                 performance = page1.find('Performance')
#                 if performance is not None:
#                     left_y = add_section_title(slide, left_x, left_y, "Performance")
#                     for table in performance.findall('Table'):
#                         left_y = add_generic_table(slide, table, left_x, left_y, LEFT_COLUMN_WIDTH)
                
#                 left_y = add_combined_section(slide, page1, left_x, left_y, LEFT_COLUMN_WIDTH)

#                 # Right Column Processing
#                 right_y = current_y
#                 execution_tracker = page1.find('ExecutionTracker/Table')
#                 if execution_tracker is not None:
#                     right_y = add_section_title(slide, right_x, right_y, "Execution Tracker Overview")
#                     right_y = add_generic_table(slide, execution_tracker, right_x, right_y, RIGHT_COLUMN_WIDTH)

#                 overdue_actions = page1.find('OverdueActions/Table')
#                 if overdue_actions is not None:
#                     right_y = add_section_title(slide, right_x, right_y, "Overdue Key Actions")
#                     right_y = add_generic_table(slide, overdue_actions, right_x, right_y, RIGHT_COLUMN_WIDTH)

#             # Process Page 2
#             page2 = result.find(".//Page[@number='2']")
#             if page2 is not None:
#                 slide = prs.slides.add_slide(blank_slide_layout)
#                 title_text = page2.find('Title').text if page2.find('Title') is not None and page2.find('Title').text else f"{result_name} - Priorities"
#                 add_slide_title(slide, title_text)
                
#                 current_y = TOP_MARGIN + Inches(0.8)
#                 priorities = page2.find('Priorities/Table')
#                 if priorities is not None:
#                     table_width = SLIDE_WIDTH - LEFT_MARGIN - RIGHT_MARGIN
#                     add_generic_table(slide, priorities, LEFT_MARGIN, current_y, table_width)

#         ppt_stream = io.BytesIO()
#         prs.save(ppt_stream)
#         ppt_stream.seek(0)
#         return ppt_stream

#     except Exception as e:
#         logging.error(f"Failed to build custom PowerPoint from XML: {e}", exc_info=True)
#         raise

# def add_slide_title(slide, text):
#     title_box = slide.shapes.add_textbox(LEFT_MARGIN, TOP_MARGIN, SLIDE_WIDTH - LEFT_MARGIN - RIGHT_MARGIN, Inches(0.4))
#     tf = title_box.text_frame
#     tf.text = text
#     tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#     tf.word_wrap = True
#     p = tf.paragraphs[0]
#     p.font.size = TITLE_FONT_SIZE
#     p.font.bold = TITLE_FONT_BOLD
#     p.font.color.rgb = TITLE_FONT_COLOR

# def add_section_title(slide, x, y, text):
#     title_box = slide.shapes.add_textbox(x, y, LEFT_COLUMN_WIDTH, Inches(0.5))
#     tf = title_box.text_frame
#     tf.text = text
#     p = tf.paragraphs[0]
#     p.font.size = SECTION_FONT_SIZE
#     p.font.bold = SECTION_FONT_BOLD
#     p.font.color.rgb = SECTION_FONT_COLOR
#     return y + Inches(0.2) + SECTION_SPACING

# def add_generic_table(slide, table_xml, x, y, width):
#     headers_xml = table_xml.find('Header')
#     cols_xml = headers_xml.findall('Column') if headers_xml is not None else table_xml.findall('Column')
#     headers = [col.text for col in cols_xml]
    
#     rows_data = []
#     for row in table_xml.findall('Row'):
#         cell_data = [cell.text if cell.text is not None else "" for cell in row.findall('Cell')]
#         if any(cell_data):
#             rows_data.append(cell_data)

#     if not rows_data: return y

#     table_shape = slide.shapes.add_table(len(rows_data) + 1, len(headers), x, y, width, MIN_ROW_HEIGHT)
#     table = table_shape.table

#     for i, header in enumerate(headers):
#         cell = table.cell(0, i)
#         cell.text = header
#         cell.fill.solid()
#         cell.fill.fore_color.rgb = TABLE_HEADER_BG_COLOR
#         p = cell.text_frame.paragraphs[0]
#         p.font.color.rgb = TABLE_HEADER_FONT_COLOR
#         p.font.size = TABLE_HEADER_FONT_SIZE
#         p.font.bold = TABLE_HEADER_FONT_BOLD

#     for r, row_data in enumerate(rows_data, 1):
#         for c, cell_data in enumerate(row_data):
#             cell = table.cell(r, c)
#             cell.text = cell_data
#             p = cell.text_frame.paragraphs[0]
#             p.font.size = TABLE_CELL_FONT_SIZE
#             p.font.color.rgb = TABLE_CELL_FONT_COLOR
    
#     autofit_table_rows(table)
#     return y + table_shape.height + SECTION_SPACING

# def add_combined_section(slide, page, x, y, width):
#     sections = {'Highlights': 'Highlights', 'Challenges': 'Challenges', 'Opportunities': 'Opportunities'}
#     has_content = False
    
#     textbox = slide.shapes.add_textbox(x, y, width, Inches(0.1))
#     tf = textbox.text_frame
#     tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#     tf.word_wrap = True

#     for section_name, section_title in sections.items():
#         section_xml = page.find(section_name)
#         if section_xml is not None:
#             key_points = [kp.text for kp in section_xml.findall('KeyPoint') if kp.text and kp.text.strip()]
#             if key_points:
#                 has_content = True
#                 p = tf.add_paragraph()
#                 p.text = section_title
#                 p.font.size = SUBSECTION_FONT_SIZE
#                 p.font.bold = SUBSECTION_FONT_BOLD
#                 p.font.color.rgb = SUBSECTION_FONT_COLOR
#                 p.space_after = SUBSECTION_SPACING
                
#                 for kp_text in key_points:
#                     p = tf.add_paragraph()
#                     p.text = "• " + kp_text
#                     p.font.size = BULLET_FONT_SIZE
#                     p.font.color.rgb = BULLET_FONT_COLOR
#                     p.space_after = BULLET_POINT_SPACING

#     if not has_content:
#         slide.shapes.element.remove(textbox.element)
#         return y
    
#     return y + textbox.height + SECTION_SPACING

# def autofit_table_rows(table):
#     col_widths = [col.width for col in table.columns]
#     for row in table.rows:
#         max_h = MIN_ROW_HEIGHT
#         for i, cell in enumerate(row.cells):
#             text_frame = cell.text_frame
#             text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#             text_frame.word_wrap = True
#             if text_frame.height > max_h:
#                 max_h = text_frame.height
#         row.height = max_h
#     for i, width in enumerate(col_widths):
#         table.columns[i].width = width
